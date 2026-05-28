from pptx import Presentation
import os

pptx_path = 'C:/Users/SEOL/Downloads/정처기앱  erd 다시 수정계획안.pptx'
output_dir = 'C:/Users/SEOL/InformationExamProject/pptx_images'
os.makedirs(output_dir, exist_ok=True)

prs = Presentation(pptx_path)

image_count = 0
for i, slide in enumerate(prs.slides, 1):
    slide_dir = os.path.join(output_dir, f'slide_{i}')
    os.makedirs(slide_dir, exist_ok=True)

    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            try:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext

                image_count += 1
                image_path = os.path.join(slide_dir, f'image_{image_count}.{image_ext}')
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                print(f"저장됨: {image_path} ({len(image_bytes)} bytes)")
            except Exception as e:
                print(f"이미지 저장 실패: {e}")

print(f"\n총 {image_count}개 이미지 저장 완료")
print(f"저장 위치: {output_dir}")
