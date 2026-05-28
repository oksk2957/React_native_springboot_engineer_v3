from pptx import Presentation
import json

pptx_path = 'C:/Users/SEOL/Downloads/정처기앱  erd 다시 수정계획안.pptx'
prs = Presentation(pptx_path)

print(f'슬라이드 수: {len(prs.slides)}')
print('=' * 50)

for i, slide in enumerate(prs.slides, 1):
    print(f'\n--- 슬라이드 {i} ---')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                print(text)
