from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

pptx_path = 'C:/Users/SEOL/Downloads/정처기앱  erd 다시 수정계획안.pptx'
prs = Presentation(pptx_path)

print(f"슬라이드 수: {len(prs.slides)}")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- 슬라이드 {i} ---")
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                # Try to decode properly
                try:
                    print(f"  텍스트: {text}")
                except:
                    print(f"  텍스트: [인코딩 문제]")

        # Check for images
        if shape.shape_type == 13:  # Picture
            print(f"  이미지 발견: {shape.name}")
            try:
                image = shape.image
                image_bytes = image.blob
                print(f"    이미지 크기: {len(image_bytes)} bytes")
                print(f"    이미지 형식: {image.ext}")
            except:
                print("    이미지 추출 실패")
